# reporting/reporter.py
# CSV 데이터를 읽어 PDF 및 PPT 보고서를 생성하는 모듈

import pandas as pd
import os
from datetime import datetime
import matplotlib.pyplot as plt
from matplotlib import font_manager, rc
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from pptx.util import Inches
from utils.logger import log_event

# --- 상수 정의 ---
LOG_DIR = 'logs'
CSV_FILE_PATH = os.path.join(LOG_DIR, 'trade_log.csv')
REPORT_DIR = 'reports'
CHART_FILE = os.path.join(REPORT_DIR, 'trade_report_chart.png')

def _setup_korean_font():
    """한글 폰트를 시스템에 맞게 설정합니다."""
    try:
        if os.name == 'nt': # Windows
            font_name = font_manager.FontProperties(fname="c:/Windows/Fonts/malgun.ttf").get_name()
            rc('font', family=font_name)
            # ReportLab 용 폰트 등록
            pdfmetrics.registerFont(TTFont('MalgunGothic', 'c:/Windows/Fonts/malgun.ttf'))
            return 'MalgunGothic'
        elif os.name == 'posix': # MacOS or Linux
            # MacOS
            if 'darwin' in os.sys.platform:
                if os.path.exists('/System/Library/Fonts/Supplemental/AppleGothic.ttf'):
                    rc('font', family='AppleGothic')
                    pdfmetrics.registerFont(TTFont('AppleGothic', '/System/Library/Fonts/Supplemental/AppleGothic.ttf'))
                    return 'AppleGothic'
            # Linux (나눔고딕이 설치되어 있다고 가정)
            elif os.path.exists('/usr/share/fonts/truetype/nanum/NanumGothic.ttf'):
                font_name = font_manager.FontProperties(fname="/usr/share/fonts/truetype/nanum/NanumGothic.ttf").get_name()
                rc('font', family=font_name)
                pdfmetrics.registerFont(TTFont('NanumGothic', '/usr/share/fonts/truetype/nanum/NanumGothic.ttf'))
                return 'NanumGothic'
        # 기본 폰트
        plt.rcParams['axes.unicode_minus'] = False # 마이너스 폰트 깨짐 방지
        return 'sans-serif' # 한글 폰트가 없을 경우 기본값
    except Exception as e:
        log_event("WARNING", f"한글 폰트 설정 중 오류 발생: {e}. 차트의 한글이 깨질 수 있습니다.")
        return 'sans-serif'

def _generate_chart(df):
    """데이터프레임을 기반으로 파이 차트를 생성하고 이미지 파일로 저장합니다."""
    if df.empty:
        return False
        
    _setup_korean_font()
    
    strategy_counts = df['strategy'].value_counts()
    
    plt.figure(figsize=(8, 6))
    plt.pie(strategy_counts, labels=strategy_counts.index, autopct='%1.1f%%', startangle=140,
            textprops={'fontsize': 14})
    plt.title('전략별 거래 비중', fontsize=16)
    plt.ylabel('') # 불필요한 y-label 제거
    plt.axis('equal') # 원형을 유지
    
    os.makedirs(REPORT_DIR, exist_ok=True)
    plt.savefig(CHART_FILE)
    plt.close() # 리소스 해제
    log_event("INFO", f"성과 분석 차트를 생성했습니다: {CHART_FILE}")
    return True

def _create_pdf_report(df, timestamp):
    """분석된 데이터와 차트를 사용하여 PDF 보고서를 생성합니다."""
    font_name = _setup_korean_font()
    report_filename = os.path.join(REPORT_DIR, f"Trade_Report_{timestamp}.pdf")
    
    c = canvas.Canvas(report_filename, pagesize=letter)
    width, height = letter
    
    # 제목
    c.setFont(font_name, 24)
    c.drawString(50, height - 70, "주간 자동매매 성과 보고서")
    
    # 생성 날짜
    c.setFont(font_name, 12)
    c.drawString(50, height - 100, f"보고일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    # 차트 이미지 삽입
    if os.path.exists(CHART_FILE):
        c.drawImage(CHART_FILE, 50, height - 400, width=400, height=300)

    # 최근 10개 거래 내역
    c.setFont(font_name, 16)
    c.drawString(50, height - 450, "최근 거래 내역")
    
    c.setFont(font_name, 8)
    y_position = height - 480
    
    # 헤더
    headers = list(df.columns)
    x_offset = 50
    for header in headers:
        c.drawString(x_offset, y_position, str(header))
        x_offset += 70

    # 데이터
    for index, row in df.tail(10).iterrows():
        y_position -= 20
        x_offset = 50
        for item in row:
            c.drawString(x_offset, y_position, str(item))
            x_offset += 70
    
    c.save()
    log_event("SUCCESS", f"PDF 보고서 생성이 완료되었습니다: {report_filename}")

def _create_ppt_report(df, timestamp):
    """분석된 데이터와 차트를 사용하여 PPT 보고서를 생성합니다."""
    report_filename = os.path.join(REPORT_DIR, f"Trade_Report_{timestamp}.pptx")
    prs = Presentation()

    # 슬라이드 1: 제목
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "주간 자동매매 성과 보고서"
    subtitle.text = f"보고일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    # 슬라이드 2: 차트와 요약
    if os.path.exists(CHART_FILE):
        chart_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(chart_slide_layout)
        title = slide.shapes.title
        title.text = "전략별 거래 비중"
        slide.shapes.add_picture(CHART_FILE, Inches(1), Inches(1.5), width=Inches(8))
    
    # 슬라이드 3: 상세 거래 내역
    table_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(table_slide_layout)
    title = slide.shapes.title
    title.text = "최근 거래 내역"
    
    rows, cols = len(df.tail(10)) + 1, len(df.columns)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
    
    # 헤더
    for col_index, col_name in enumerate(df.columns):
        table.cell(0, col_index).text = col_name

    # 데이터
    for row_index, row_data in enumerate(df.tail(10).values):
        for col_index, cell_data in enumerate(row_data):
            table.cell(row_index + 1, col_index).text = str(cell_data)

    prs.save(report_filename)
    log_event("SUCCESS", f"PPT 보고서 생성이 완료되었습니다: {report_filename}")

def generate_reports():
    """메인 보고서 생성 함수"""
    log_event("INFO", "성과 보고서 생성을 시작합니다...")
    
    if not os.path.exists(CSV_FILE_PATH):
        log_event("WARNING", f"거래 로그 파일({CSV_FILE_PATH})이 없어 보고서를 생성할 수 없습니다.")
        return

    try:
        df = pd.read_csv(CSV_FILE_PATH)
        if df.empty:
            log_event("WARNING", "거래 내역이 없어 보고서를 생성하지 않습니다.")
            return
            
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        # 1. 차트 생성
        if _generate_chart(df):
            # 2. PDF 생성
            _create_pdf_report(df, timestamp)
            # 3. PPT 생성
            _create_ppt_report(df, timestamp)
            
            # 생성된 차트 파일 삭제 (보고서에 포함되었으므로)
            os.remove(CHART_FILE)
            
    except Exception as e:
        log_event("ERROR", f"보고서 생성 중 심각한 오류 발생: {e}")

if __name__ == '__main__':
    # 테스트용 코드
    # 임시 trade_log.csv 생성
    os.makedirs(LOG_DIR, exist_ok=True)
    dummy_data = {
        'timestamp': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')] * 2,
        'strategy': ['long_term', 'short_term'],
        'symbol': ['005930', '000660'],
        'action': ['BUY', 'SELL'],
        'price': ['70,000', '120,000'],
        'quantity': [10, 5],
        'total_amount': ['700,000', '600,000'],
        'reason': ['테스트 매수', '테스트 매도']
    }
    dummy_df = pd.DataFrame(dummy_data)
    dummy_df.to_csv(CSV_FILE_PATH, index=False, encoding='utf-8-sig')

    generate_reports() 